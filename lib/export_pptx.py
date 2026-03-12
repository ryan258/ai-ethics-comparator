"""
PowerPoint (PPTX) export (Phase 6).

Generates a presentation mirroring the PDF report sections.
Requires ``python-pptx`` (optional dependency).
"""

from __future__ import annotations

import io
import logging
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ModuleNotFoundError:  # pragma: no cover
    Presentation = None  # type: ignore[assignment,misc]


def pptx_available() -> bool:
    return Presentation is not None


def generate_pptx(
    run_data: Dict[str, Any],
    paradox: Dict[str, Any],
    insight: Optional[Dict[str, Any]] = None,
) -> bytes:
    """Generate PPTX bytes for a single run."""
    if Presentation is None:
        raise RuntimeError(
            "python-pptx is required for PowerPoint export.  "
            "Install it with: pip install python-pptx"
        )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, run_data, paradox)
    _add_distribution_slide(prs, run_data, paradox)
    if isinstance(insight, dict) and insight.get("content"):
        _add_analysis_slide(prs, insight)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Helpers ──────────────────────────────────────────────────────────────────

_BG = RGBColor(0x12, 0x12, 0x12) if Presentation else None
_TEXT = RGBColor(0xEB, 0xD2, 0xBE) if Presentation else None
_ACCENT = RGBColor(0xA6, 0xAC, 0xCD) if Presentation else None
_SUCCESS = RGBColor(0x98, 0xC3, 0x79) if Presentation else None


def _set_slide_bg(slide: Any, color: Any) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: int = 18,
    color: Any = None,
    bold: bool = False,
    alignment: Any = None,
) -> Any:
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    if color:
        p.font.color.rgb = color
    p.font.bold = bold
    if alignment:
        p.alignment = alignment
    return txBox


def _add_title_slide(prs: Any, run_data: Dict[str, Any], paradox: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, _BG)

    _add_textbox(slide, 0.8, 0.5, 11, 0.5,
                 "AI ETHICS COMPARATOR  —  EXECUTIVE BRIEF",
                 font_size=12, color=_ACCENT, bold=True)

    _add_textbox(slide, 0.8, 1.2, 11, 1.2,
                 "Ethical Decision Report",
                 font_size=36, color=_TEXT, bold=True)

    _add_textbox(slide, 0.8, 2.6, 11, 0.8,
                 paradox.get("title", "Unknown paradox"),
                 font_size=22, color=_SUCCESS, bold=True)

    # Key finding
    summary = run_data.get("summary", {})
    options = run_data.get("options", [])
    option_lookup = {o["id"]: o.get("label", f"Option {o['id']}") for o in options if isinstance(o, dict)}
    summary_opts = summary.get("options", []) if isinstance(summary, dict) else []
    max_count = max((int(o.get("count", 0) or 0) for o in summary_opts if isinstance(o, dict)), default=0)
    leaders = [option_lookup.get(o.get("id"), "?") for o in summary_opts
               if isinstance(o, dict) and int(o.get("count", 0) or 0) == max_count and max_count]
    lead_label = " / ".join(leaders) if leaders else "No dominant choice"

    _add_textbox(slide, 0.8, 3.8, 11, 0.4,
                 "KEY FINDING", font_size=11, color=_ACCENT, bold=True)
    _add_textbox(slide, 0.8, 4.3, 11, 0.8,
                 lead_label, font_size=28, color=_TEXT, bold=True)

    total = sum(int(o.get("count", 0) or 0) for o in summary_opts if isinstance(o, dict))
    support = f"{max_count} of {total} responses" if total else ""
    _add_textbox(slide, 0.8, 5.2, 11, 0.4, support, font_size=14, color=_SUCCESS, bold=True)

    # Metadata
    meta = f"Model: {run_data.get('modelName', '?')}  |  Run: {run_data.get('runId', '?')}"
    _add_textbox(slide, 0.8, 6.5, 11, 0.4, meta, font_size=10, color=_ACCENT)


def _add_distribution_slide(prs: Any, run_data: Dict[str, Any], paradox: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _BG)

    _add_textbox(slide, 0.8, 0.5, 11, 0.5,
                 "Choice Distribution", font_size=24, color=_TEXT, bold=True)

    summary = run_data.get("summary", {})
    options = run_data.get("options", [])
    option_lookup = {o["id"]: o for o in options if isinstance(o, dict)}
    summary_opts = summary.get("options", []) if isinstance(summary, dict) else []

    y = 1.4
    for opt in summary_opts:
        if not isinstance(opt, dict):
            continue
        oid = opt.get("id")
        meta = option_lookup.get(oid, {})
        label = meta.get("label", f"Option {oid}")
        count = int(opt.get("count", 0) or 0)
        pct = float(opt.get("percentage", 0.0) or 0.0)

        _add_textbox(slide, 0.8, y, 8, 0.35,
                     f"{label}:  {count}  ({pct:.1f}%)",
                     font_size=14, color=_TEXT, bold=True)
        y += 0.5

        desc = meta.get("description", "")
        if desc:
            _add_textbox(slide, 0.8, y, 10, 0.3, desc, font_size=10, color=_ACCENT)
            y += 0.45


def _add_analysis_slide(prs: Any, insight: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _BG)

    _add_textbox(slide, 0.8, 0.5, 11, 0.5,
                 "Analyst Assessment", font_size=24, color=_TEXT, bold=True)

    content = insight.get("content", {})
    fw = content.get("dominant_framework", "")
    if fw:
        _add_textbox(slide, 0.8, 1.3, 11, 0.4,
                     "DOMINANT FRAMEWORK", font_size=11, color=_ACCENT, bold=True)
        _add_textbox(slide, 0.8, 1.8, 11, 0.6, fw, font_size=22, color=_SUCCESS, bold=True)

    insights_list = content.get("key_insights", [])
    y = 2.8
    for ki in insights_list[:5]:
        _add_textbox(slide, 0.8, y, 11, 0.35, f"  •  {ki}", font_size=12, color=_TEXT)
        y += 0.45
